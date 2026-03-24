from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any
import base64
import mimetypes

from docx import Document
from moviepy import VideoFileClip
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation

from models import ParseResult


VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"}

def _guess_image_mime(image_bytes: bytes, filename: str | None = None) -> str:
    if filename:
        guessed, _ = mimetypes.guess_type(filename)
        if guessed and guessed.startswith("image/"):
            return guessed

    if image_bytes.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if image_bytes.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if image_bytes.startswith((b"GIF87a", b"GIF89a")):
        return "image/gif"
    if image_bytes.startswith(b"RIFF") and image_bytes[8:12] == b"WEBP":
        return "image/webp"
    if image_bytes.startswith((b"II*\x00", b"MM\x00*")):
        return "image/tiff"
    return "image/png"

def _describe_images_with_openai(images: list[tuple[bytes, str | None]], client: OpenAI | None) -> list[str]:
    if not images:
        return []

    descriptions: list[str] = []
    for i, (image_bytes, filename) in enumerate(images, start=1):
        if len(image_bytes) > 5_000_000:
            descriptions.append(f"[Image {i}] Embedded image omitted (file too large to process).")
            continue

        if client is None:
            descriptions.append(f"[Image {i}] Embedded image extracted from file.")
            continue

        mime_type = _guess_image_mime(image_bytes, filename=filename)
        encoded = base64.b64encode(image_bytes).decode("utf-8")
        data_url = f"data:{mime_type};base64,{encoded}"

        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this embedded assignment image for grading context in 1-2 concise sentences.",
                            },
                            {"type": "image_url", "image_url": {"url": data_url}},
                        ],
                    }
                ],
                max_tokens=140,
            )
            content = response.choices[0].message.content or "No description returned."
            descriptions.append(f"[Image {i}] {content.strip()}")
        except Exception:
            descriptions.append(f"[Image {i}] Embedded image extracted from file.")

    return descriptions

def parse_docx(path: Path, openai_client: OpenAI | None = None) -> tuple[str, dict[str, Any]]:
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    image_bytes = [
        (
            rel.target_part.blob,
            str(getattr(getattr(rel, "target_part", None), "partname", "")) or None,
        )
        for rel in doc.part.rels.values()
        if "image" in rel.reltype and hasattr(rel, "target_part") and hasattr(rel.target_part, "blob")
    ]
    image_descriptions = _describe_images_with_openai(image_bytes, openai_client)
    combined_text = text
    if image_descriptions:
        combined_text = f"{text}\n\nEmbedded Images:\n" + "\n".join(image_descriptions)
    return combined_text.strip(), {
        "paragraph_count": len(doc.paragraphs),
        "image_count": len(image_bytes),
        "images_described": len(image_descriptions),
    }


def parse_pdf(path: Path, openai_client: OpenAI | None = None) -> tuple[str, dict[str, Any]]:
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    image_bytes: list[tuple[bytes, str | None]] = []
    for page in reader.pages:
        for image in getattr(page, "images", []):
            data = getattr(image, "data", None)
            if data:
                image_bytes.append((data, getattr(image, "name", None)))
    image_descriptions = _describe_images_with_openai(image_bytes, openai_client)
    text = "\n".join(pages).strip()
    if image_descriptions:
        text = f"{text}\n\nEmbedded Images:\n" + "\n".join(image_descriptions)
    return text.strip(), {
        "page_count": len(reader.pages),
        "image_count": len(image_bytes),
        "images_described": len(image_descriptions),
    }


def parse_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    presentation = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("[Speaker Notes]")
                lines.append(notes)
    return "\n".join(lines), {"slide_count": len(presentation.slides)}


def parse_video(path: Path, client: OpenAI) -> tuple[str, dict[str, Any]]:
    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_audio:
        temp_audio_path = Path(temp_audio.name)

    try:
        clip = VideoFileClip(str(path))
        duration = clip.duration
        clip.audio.write_audiofile(str(temp_audio_path), logger=None)
        clip.close()

        with temp_audio_path.open("rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="gpt-4o-mini-transcribe",
                file=audio_file,
            )

        return transcript.text.strip(), {"duration_seconds": duration}
    finally:
        if temp_audio_path.exists():
            os.remove(temp_audio_path)


def detect_assignment_type(file_extension: str) -> str:
    ext = file_extension.lower()
    if ext in {".pdf", ".doc", ".docx", ".txt", ".rtf"}:
        return "Written Assignment"
    if ext in {".ppt", ".pptx"}:
        return "PowerPoint"
    if ext in VIDEO_EXTENSIONS:
        return "Video"
    raise ValueError(f"Unsupported submission type: {ext}")


def parse_submission_file(path: Path, openai_client: OpenAI) -> ParseResult:
    ext = path.suffix.lower()
    assignment_type = detect_assignment_type(ext)

    if ext in {".docx"}:
        text, metadata = parse_docx(path, openai_client=openai_client)
    elif ext in {".pdf"}:
        text, metadata = parse_pdf(path, openai_client=openai_client)
    elif ext in {".pptx"}:
        text, metadata = parse_pptx(path)
    elif ext in VIDEO_EXTENSIONS:
        text, metadata = parse_video(path, openai_client)
    elif ext in {".txt", ".rtf", ".doc", ".ppt"}:
        text = path.read_text(encoding="utf-8", errors="ignore")
        metadata = {"fallback_parser": True}
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

    if not text.strip():
        raise ValueError("No extractable content found in submission")

    return ParseResult(
        assignment_type=assignment_type,
        extracted_text=text,
        parser_metadata=metadata,
    )
